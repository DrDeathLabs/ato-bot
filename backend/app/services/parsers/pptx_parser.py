"""PowerPoint parser — extract slide text via python-pptx."""
from __future__ import annotations

from pathlib import Path

from app.services.parsers.base import ParsedDocument, ParsedPage


def parse_pptx(file_path: str) -> ParsedDocument:
    """Parse a .pptx/.ppt file and return one ParsedPage per slide."""
    doc = ParsedDocument(filename=Path(file_path).name, file_type="pptx")
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        pages: list[ParsedPage] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            title: str | None = None

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
                # Identify slide title placeholder (idx == 0)
                # Must check shape.is_placeholder first — accessing placeholder_format
                # on a non-placeholder shape raises "shape is not a placeholder"
                try:
                    if shape.is_placeholder and shape.placeholder_format.idx == 0:
                        title = shape.text_frame.text.strip()
                except Exception:
                    pass

            if texts:
                pages.append(ParsedPage(
                    page_number=slide_idx,
                    content="\n".join(texts),
                    section_title=title or f"Slide {slide_idx}",
                ))

        doc.pages = pages
        doc.full_text = "\n\n".join(p.content for p in pages)

    except Exception as e:
        doc.error = str(e)

    return doc

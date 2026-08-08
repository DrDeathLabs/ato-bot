"""Executive PowerPoint briefing report."""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy import select

from app.core.config import get_settings
from app.core.database import AsyncSessionLocal
from app.models.orm import Assessment, ControlFinding, Project

settings = get_settings()

BLUE = RGBColor(0x1F, 0x4E, 0x79)
GREEN = RGBColor(0x00, 0x70, 0x00)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xCC, 0x00, 0x00)
GRAY = RGBColor(0x80, 0x80, 0x80)


async def generate_pptx(assessment_id: int) -> str:
    async with AsyncSessionLocal() as db:
        result = await db.execute(select(Assessment).where(Assessment.id == assessment_id))
        assessment = result.scalar_one()
        proj_result = await db.execute(select(Project).where(Project.id == assessment.project_id))
        project = proj_result.scalar_one()
        findings_result = await db.execute(
            select(ControlFinding)
            .where(ControlFinding.assessment_id == assessment_id)
            .order_by(ControlFinding.control_family, ControlFinding.control_id)
        )
        findings = findings_result.scalars().all()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    # ── Title Slide ────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_text(slide, "Security Assessment Executive Briefing", Inches(1), Inches(2), Inches(11), Inches(1.5),
               size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    _add_text(slide, project.name, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
               size=24, color=BLUE, align=PP_ALIGN.CENTER)
    _add_text(slide, f"NIST 800-53 Rev 5 {project.impact_baseline.upper()} Baseline",
               Inches(1), Inches(4.3), Inches(11), Inches(0.6), size=18, align=PP_ALIGN.CENTER)

    # ── Compliance Summary ─────────────────────────────────────────────────────
    total = len(findings)
    compliant = sum(1 for f in findings if f.status == "compliant")
    partial = sum(1 for f in findings if f.status == "partially_compliant")
    non = sum(1 for f in findings if f.status == "non_compliant")
    na = sum(1 for f in findings if f.status == "not_applicable")
    score = round((compliant + na) / max(total, 1) * 100, 1)

    slide2 = prs.slides.add_slide(blank_layout)
    _add_text(slide2, "Compliance Summary", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
               size=28, bold=True, color=BLUE)

    stats = [
        (f"{score}%", "Overall Score", BLUE),
        (f"{compliant}", "Compliant", GREEN),
        (f"{partial}", "Partial", ORANGE),
        (f"{non}", "Non-Compliant", RED),
        (f"{na}", "N/A", GRAY),
    ]
    for i, (num, label, color) in enumerate(stats):
        x = Inches(0.5 + i * 2.5)
        _add_text(slide2, num, x, Inches(1.5), Inches(2.2), Inches(1.2), size=40, bold=True, color=color, align=PP_ALIGN.CENTER)
        _add_text(slide2, label, x, Inches(2.7), Inches(2.2), Inches(0.5), size=14, color=color, align=PP_ALIGN.CENTER)

    # ── Results by Family ──────────────────────────────────────────────────────
    families = sorted({f.control_family for f in findings})
    family_rows = []
    for fam in families:
        fam_f = [f for f in findings if f.control_family == fam]
        c = sum(1 for f in fam_f if f.status == "compliant")
        p = sum(1 for f in fam_f if f.status == "partially_compliant")
        n = sum(1 for f in fam_f if f.status == "non_compliant")
        fam_score = round(c / len(fam_f) * 100, 0)
        family_rows.append((fam, len(fam_f), c, p, n, f"{fam_score}%"))

    slide3 = prs.slides.add_slide(blank_layout)
    _add_text(slide3, "Results by Control Family", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
               size=28, bold=True, color=BLUE)
    headers = ["Family", "Total", "Compliant", "Partial", "Non-Compliant", "Score"]
    _add_table(slide3, headers, family_rows, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))

    # ── Top Gaps ───────────────────────────────────────────────────────────────
    slide4 = prs.slides.add_slide(blank_layout)
    _add_text(slide4, "Top Priority Findings (Non-Compliant)", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
               size=28, bold=True, color=BLUE)
    non_findings = [f for f in findings if f.status == "non_compliant"][:10]
    y = Inches(1.3)
    for f in non_findings:
        _add_text(slide4, f"{f.control_id}: {f.control_title}", Inches(0.5), y, Inches(12), Inches(0.4),
                   size=12, bold=True, color=RED)
        if f.gaps:
            _add_text(slide4, f.gaps[0][:120], Inches(0.8), y + Inches(0.4), Inches(11.5), Inches(0.35),
                       size=10, color=GRAY)
        y += Inches(0.8)
        if y > Inches(6.5):
            break

    output_path = Path(settings.output_dir) / f"assessment_{assessment_id}_executive.pptx"
    os.makedirs(settings.output_dir, exist_ok=True)
    prs.save(str(output_path))
    return str(output_path)


def _add_text(slide, text: str, left, top, width, height,
              size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _add_table(slide, headers, rows, left, top, width, height):
    cols = len(headers)
    table = slide.shapes.add_table(len(rows) + 1, cols, left, top, width, height).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = str(val)
